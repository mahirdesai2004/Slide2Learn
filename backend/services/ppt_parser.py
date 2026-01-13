from pptx import Presentation

def extract_slides_text(file_path: str):
    prs = Presentation(file_path)
    slides_data = []

    for idx, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    text_parts.append(text)

        slide_text = "\n".join(text_parts)

        slides_data.append({
            "slide_no": idx,
            "raw_text": slide_text
        })

    return slides_data